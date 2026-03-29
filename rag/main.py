import os
import io
from pathlib import Path
import sys
from urllib import request
from dotenv import load_dotenv
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_community.vectorstores import Chroma
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.prompts import ChatPromptTemplate

import pdfplumber
from pptx import Presentation
from typing import List
from fastapi import FastAPI, Form, File, UploadFile, Body, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

# load environment variables from .env.local
# __file__ is the current file. .parent gets the folder
# .parent.parent moves up one level
env_path = Path(__file__).resolve().parent.parent / '.env.local'

# load the specific file
load_dotenv(dotenv_path=env_path)


# get variable
openai_api_key = os.getenv("OPENAI_API_KEY")

if not openai_api_key:
    print("OPENAI_API_KEY is not accessible from this file.")
    sys.exit(1)


app = FastAPI()

origins = [
    "http://localhost:3000",
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

model = ChatOpenAI(model_name="gpt-4", temperature=0.7)


@app.get("/")
def root():
    return {"Hello" : "world"}

# initialize OpenAI embeddings
embeddings = OpenAIEmbeddings()

# create chunking function using RecursiveCharacterTextSplitter
text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=1000,
    chunk_overlap=200
)

# create the vector database using Chroma from LangChain
vector_store = Chroma(
    embedding_function=embeddings,
    persist_directory="./chroma_db"
)

# extract text from a pdf
def extract_text_from_pdf(file_bytes):
    text_content = ""
    # Wrap bytes in BytesIO so that it looks like a file to the library
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                text_content += text + "\n"

    return text_content


# extract text from a pptx
def extract_text_from_pptx(file_bytes):
    text_content = ""
    # laod presentation from bytes
    presentation = Presentation(io.BytesIO(file_bytes))

    for slide in presentation.slides:
        # Iterate through all "shapes" (text boxes, titles, etc.) on the slide
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content += shape.text + "\n"

    return text_content



@app.post("/api/create-study-buddy")
async def create_study_buddy(
    title: str = Form(...), 
    description: str = Form(...), 
    category: str = Form(...), 
    files: List[UploadFile] = File(...),
    buddyId: str = Form(...)
):
    print(f"The received buddyId is {buddyId}")
    total_chunks = 0

    status = "success"
    
    # chunk each file instead of the entire text to preserve context
    for file in files:
        try:
            file_bytes = await file.read()

            filename = file.filename.lower()

            file_content = ""

            if filename.endswith(".pdf"):
                print(f"{filename} is a pdf file")
                file_content = extract_text_from_pdf(file_bytes)

            if filename.endswith(".pptx"):
                print(f"{filename} is a pptx file")
                file_content = extract_text_from_pptx(file_bytes)

            
            if not file_content:
                continue

            # split file_content into chunks
            chunks = text_splitter.split_text(file_content)

            # create metadata
            metadata = [{"buddy_id": buddyId, "chunk_index": i} for i, _ in enumerate(chunks)] # make sure that the length of chunks is the same as the length of the metadata
            

            # add to the vector database
            if chunks:
                print(f"The metadata attached to each chunk is {metadata[0]}")
                vector_store.add_texts(texts=chunks, metadatas=metadata)
                total_chunks += len(chunks)

        
        except Exception as e:
            print(f"Error processing {file.filename}", e)
            status = "error"

    


    return {
        "status": status, 
        "buddy_id": buddyId, 
        "chunks_count": total_chunks
    }


# Chat Request Template
class ChatRequest(BaseModel):
    question: str
    buddy_id: str # To find chunks that correspond to the study bot


# chat request endpoint
@app.post("/api/chat")
async def chat_request(request: ChatRequest):
    print(f"Question: {request.question} with Buddy with ID {request.buddy_id}")

    # convert the query to an embedding and search for the 3 most relevant chunks
    query = request.question
    results = vector_store.similarity_search(query=query, k=3, filter={"buddy_id": request.buddy_id})

    if not results:
        return {"error": "Could not find any chunks related to your question"}
    
    
    # combine the retrieved chunks into 1 string
    context_text = "\n\n---\n\n".join([doc.page_content for doc in results])


    # Engineer a prompt
    system_prompt = """
    You are an AI assistant that is a study buddy. You use your pre-existing knowledge as well as context provided
    to you to accurately answer queries made by a user. The context you receive will come from documents uploaded by the
    user related to the topic that they want the study buddy to help them with. If the user's query is not plausible or
    is inappropriate, say that you are unable to help.

    If the user's query is a simple question like a 'yes or no' or other type of simple question, then a simple
    response will be sufficient. However, most of the user's querys will be asking what something is and explaining
    it to them. So it is important to be very detailed in your answer and give the user a good understanding
    of the query.

    ------
    Context: {context_text}
    ---
    Query: {query}
    """

    prompt_template = ChatPromptTemplate.from_template(system_prompt)
    prompt = prompt_template.format(context_text=context_text, query=request.question)

    # Generate the answer using OpenAI
    response_text = model.invoke(prompt)

    # save the response
    response = response_text.content

    return {
        "response": response,
    }
    




# temp
@app.post("/api/test-endpoint")
def testing_post(text: str = Body(...), number: int = Body(...)):
    print(f"The text received is:\n{text}\nWith number: {number}")
    return True